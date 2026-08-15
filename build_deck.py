# MediCluster Definitive Deck Generator
# Consolidates content from MediCluster_Full_Report.md + all 3 existing decks
# Design language: indigo/violet/teal, white cards, slate body, 13.33x7.5 widescreen
# Same design DNA as the existing gen_pptx.py (same RGB colors and helpers)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import datetime

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── DESIGN SYSTEM (matches gen_pptx.py) ─────────────────────────────────────
INDIGO    = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_D  = RGBColor(0x37, 0x30, 0xA3)
VIOLET    = RGBColor(0x7C, 0x3A, 0xED)
SLATE     = RGBColor(0x33, 0x41, 0x55)
MUTED     = RGBColor(0x64, 0x74, 0x8B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xEE, 0xF2, 0xFF)
RED       = RGBColor(0xDC, 0x26, 0x26)
GREEN     = RGBColor(0x16, 0xA3, 0x4A)
AMBER     = RGBColor(0xD9, 0x77, 0x06)
TEAL      = RGBColor(0x05, 0x96, 0x69)
BLUE      = RGBColor(0x25, 0x63, 0xEB)
SKY       = RGBColor(0x03, 0x82, 0xF6)
BG        = RGBColor(0xF8, 0xFA, 0xFC)
BORDER    = RGBColor(0xE2, 0xE8, 0xF0)
RISK_LOW  = RGBColor(0x16, 0xA3, 0x4A)
RISK_MOD  = RGBColor(0xD9, 0x77, 0x06)
RISK_HIGH = RGBColor(0xEA, 0x58, 0x0C)
RISK_CRIT = RGBColor(0xDC, 0x26, 0x26)

BLANK = prs.slide_layouts[6]

# ── PRIMITIVES ──────────────────────────────────────────────────────────────
def fill(shape, color):
    f = shape.fill
    f.solid()
    f.fore_color.rgb = color

def add_rect(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shp, color)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    return shp

def add_round(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shp, color)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp

def add_text(slide, text, x, y, w, h, size=14, bold=False, color=SLATE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tb

def add_multi(slide, runs, x, y, w, h, align=PP_ALIGN.LEFT):
    """runs: list of (text, size, bold, color) — first run is on first paragraph"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    for i, (txt, sz, bld, clr) in enumerate(runs):
        if "\n" in txt:
            parts = txt.split("\n")
            for j, part in enumerate(parts):
                if j > 0:
                    p = tf.add_paragraph()
                    p.alignment = align
                r = p.add_run()
                r.text = part
                r.font.size = Pt(sz)
                r.font.bold = bld
                r.font.color.rgb = clr
                r.font.name = "Calibri"
        else:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(sz)
            r.font.bold = bld
            r.font.color.rgb = clr
            r.font.name = "Calibri"
    return tb

def bg(slide, color=BG):
    b = slide.background
    b.fill.solid()
    b.fill.fore_color.rgb = color

def header(slide, title, sub=None, accent=INDIGO):
    add_rect(slide, 0, 0, 13.333, 1.05, accent)
    add_rect(slide, 0, 1.05, 13.333, 0.06, VIOLET)
    add_text(slide, title, 0.4, 0.12, 11.5, 0.55, size=24, bold=True, color=WHITE)
    if sub:
        add_text(slide, sub, 0.4, 0.62, 11.5, 0.38, size=12, color=LIGHT, italic=True)
    # page brand
    add_text(slide, "MediCluster  ·  Patient Health Risk Segregation",
             0.4, 7.15, 9, 0.3, size=9, color=MUTED, italic=True)
    return 1.2  # content Y start

def card(slide, x, y, w, h, title, lines, color=INDIGO, body=SLATE, body_size=11, badge=None):
    add_rect(slide, x, y, w, h, WHITE, line=BORDER)
    add_rect(slide, x, y, w, 0.42, color)
    add_text(slide, title, x+0.12, y+0.05, w-0.24, 0.32, size=12, bold=True, color=WHITE)
    if badge:
        add_rect(slide, x+w-0.85, y+0.07, 0.75, 0.28, AMBER)
        add_text(slide, badge, x+w-0.85, y+0.09, 0.75, 0.26, size=8, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
    body_y = y + 0.5
    body_h = h - 0.6
    tb = slide.shapes.add_textbox(Inches(x+0.12), Inches(body_y), Inches(w-0.24), Inches(body_h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = u"•  " + line
        r.font.size = Pt(body_size)
        r.font.color.rgb = body
        r.font.name = "Calibri"

def bullet_block(slide, items, x, y, w, h, size=12, color=SLATE, title=None, title_color=INDIGO):
    cur_y = y
    if title:
        add_text(slide, title, x, cur_y, w, 0.36, size=14, bold=True, color=title_color)
        cur_y += 0.4
        h -= 0.4
    tb = slide.shapes.add_textbox(Inches(x), Inches(cur_y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = u"•  " + item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"

def table(slide, x, y, w, h, headers, rows, header_color=INDIGO,
          row_alt=BG, header_text_color=WHITE, row_text_color=SLATE,
          col_widths=None, font_size=11, header_size=12):
    """Manual table built from rects + textboxes (no python-pptx native table)."""
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [w / n_cols] * n_cols
    else:
        total = sum(col_widths)
        col_widths = [cw * w / total for cw in col_widths]
    # header
    cx = x
    for i, hd in enumerate(headers):
        add_rect(slide, cx, y, col_widths[i], 0.4, header_color)
        add_text(slide, hd, cx+0.05, y+0.04, col_widths[i]-0.1, 0.32,
                 size=header_size, bold=True, color=header_text_color, align=PP_ALIGN.LEFT)
        cx += col_widths[i]
    # rows
    row_h = (h - 0.4) / max(1, len(rows))
    for ri, row in enumerate(rows):
        ry = y + 0.4 + ri * row_h
        bg_color = WHITE if ri % 2 == 0 else row_alt
        add_rect(slide, x, ry, w, row_h, bg_color, line=BORDER)
        cx = x
        for i, val in enumerate(row):
            add_text(slide, str(val), cx+0.05, ry+0.04, col_widths[i]-0.1, row_h-0.08,
                     size=font_size, color=row_text_color, align=PP_ALIGN.LEFT)
            cx += col_widths[i]

def stat(slide, x, y, w, h, big, small, color=INDIGO, bg_color=INDIGO_D):
    add_rect(slide, x, y, w, h, bg_color)
    add_text(slide, big, x, y+0.15, w, 0.7, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, small, x, y+0.9, w, 0.6, size=10, color=LIGHT, align=PP_ALIGN.CENTER)

def title_slide():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, 13.333, 7.5, INDIGO)
    add_rect(s, 0, 4.6, 13.333, 2.9, VIOLET)
    add_rect(s, 0, 4.55, 13.333, 0.05, RGBColor(0xA7, 0x8B, 0xFA))
    add_text(s, "MediCluster", 1, 1.2, 11.333, 1.5, size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Patient Health Risk Segregation Platform",
             1, 2.85, 11.333, 0.7, size=24, color=LIGHT, align=PP_ALIGN.CENTER)
    add_text(s, "AI-Powered Clinical Decision Support  ·  Clustering · NLP · Imaging · Forecasting · Dispatch",
             1, 3.55, 11.333, 0.5, size=13, color=RGBColor(0xC7, 0xD2, 0xFE), align=PP_ALIGN.CENTER, italic=True)
    # chips for the 4 product pillars
    for i, label in enumerate(["ML Clustering", "LungLens Imaging", "CliniQ NLP", "ARIA Dispatch"]):
        cx = 0.8 + i * 3.0
        add_rect(s, cx, 4.95, 2.7, 1.55, INDIGO_D)
        add_text(s, label, cx, 5.35, 2.7, 0.8, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Comprehensive Technical Overview  ·  " + datetime.date.today().strftime("%B %Y"),
             0, 6.75, 13.333, 0.4, size=11, color=RGBColor(0xC7, 0xD2, 0xFE), align=PP_ALIGN.CENTER, italic=True)
    add_text(s, "CtrlAltElite  ·  M S Engineering College  ·  May 2026",
             0, 7.1, 13.333, 0.35, size=10, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.CENTER, italic=True)

def agenda():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Agenda", "What this deck covers — 30 slides across 9 sections")
    cols = [
        ("01", "Foundation",     ["Problem & Goals", "Architecture", "Tech Stack", "Data Flow"]),
        ("02", "ML Engine",      ["Clustering Algos", "Risk Scoring", "SHAP", "Anomalies", "AutoML"]),
        ("03", "Clinical AI",    ["NLP Notes", "Vitals Forecast", "MEWS", "Imaging AI", "Drug Check"]),
        ("04", "Frontend & API", ["13 Pages", "API Reference", "Real-time UX"]),
        ("05", "Deployment",     ["Docker", "Performance", "Security", "Testing"]),
        ("06", "Roadmap",        ["Phase 1-3 Plans", "Limitations", "Glossary"]),
    ]
    for i, (num, title, items) in enumerate(cols):
        col = i % 3
        row = i // 3
        cx = 0.4 + col * 4.32
        cy = y0 + 0.2 + row * 2.85
        add_rect(s, cx, cy, 4.1, 2.6, WHITE, line=BORDER)
        add_rect(s, cx, cy, 4.1, 0.55, INDIGO)
        add_text(s, num, cx+0.15, cy+0.08, 0.6, 0.4, size=18, bold=True, color=WHITE)
        add_text(s, title, cx+0.85, cy+0.1, 3.1, 0.4, size=15, bold=True, color=WHITE)
        add_text(s, "\n".join(["• " + i_ for i_ in items]),
                 cx+0.18, cy+0.7, 3.8, 1.85, size=12, color=SLATE)

# ═════════════════════════════════════════════════════════════════════════════
# SECTION 1 — PROBLEM & GOALS
# ═════════════════════════════════════════════════════════════════════════════
def slide_problem():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "The Problem", "Why patient risk segregation needs a new approach")
    add_text(s, "Modern healthcare is overwhelmed. Emergency departments face hundreds of patients daily with limited staff, fragmented data, and rising acuity. MediCluster addresses five concrete failure modes.",
             0.4, y0, 12.5, 0.8, size=12, color=MUTED, italic=True)
    problems = [
        ("Triage bottlenecks",  RED,   "Manual assessment of hundreds of patients daily"),
        ("Missed deterioration",AMBER, "Vital signs not continuously monitored between rounds"),
        ("Unstructured data",   VIOLET,"Doctor notes locked in free text — no ICD-10 coding"),
        ("Imaging backlogs",    BLUE,  "Radiologists overwhelmed with X-ray and CT queues"),
        ("Medication errors",   TEAL,  "Drug interactions missed during prescription"),
        ("No labelled data",    INDIGO,"Supervised risk models impossible without ground truth"),
    ]
    for i, (title, clr, desc) in enumerate(problems):
        col = i % 3
        row = i // 3
        cx = 0.4 + col * 4.3
        cy = y0 + 0.95 + row * 2.55
        card(s, cx, cy, 4.1, 2.35, title, [desc], color=clr, body=SLATE, body_size=12)

def slide_goals_users():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Design Goals & Target Users", "Five principles that shaped every technical decision")
    add_text(s, "DESIGN GOALS", 0.4, y0, 6.0, 0.4, size=14, bold=True, color=INDIGO)
    add_text(s, "TARGET USERS", 6.8, y0, 6.0, 0.4, size=14, bold=True, color=VIOLET)
    goals = [
        "Zero labelled data required — core clustering works on raw patient vitals",
        "Graceful degradation — every ML component falls back to rule-based logic if libraries are absent",
        "Clinical interpretability — every prediction includes plain-language explanations for clinicians",
        "Real-time operation — live vital monitoring, instant NLP analysis, real-time dispatch",
        "Full observability — SHAP, MEWS, data quality flags, and confidence intervals on every output",
    ]
    bullet_block(s, goals, 0.4, y0+0.45, 6.0, 5.0, size=12, color=SLATE)
    users = [
        ("Emergency Physicians", "Real-time triage support, critical alerts"),
        ("Ward Nurses",          "Continuous vital monitoring, MEWS early warning"),
        ("Radiologists",         "AI-assisted chest X-ray preliminary reading"),
        ("Administrators",       "Population-level risk dashboard, workload forecast"),
        ("Dispatchers",          "GPS-integrated real-time ambulance coordination"),
    ]
    for i, (t, d) in enumerate(users):
        cy = y0 + 0.5 + i * 0.95
        add_round(s, 6.8, cy, 6.1, 0.85, WHITE, line=BORDER)
        add_rect(s, 6.8, cy, 0.15, 0.85, VIOLET)
        add_text(s, t, 7.0, cy+0.08, 6.0, 0.32, size=12, bold=True, color=VIOLET)
        add_text(s, d, 7.0, cy+0.42, 6.0, 0.38, size=11, color=SLATE)

# ═════════════════════════════════════════════════════════════════════════════
# SECTION 2 — ARCHITECTURE & STACK
# ═════════════════════════════════════════════════════════════════════════════
def slide_architecture():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "System Architecture", "Three-tier microservice with persistent storage")
    layers = [
        (INDIGO, "Client Tier",  "React 18  ·  Vite SPA  ·  TailwindCSS  ·  Recharts  ·  D3.js  ·  Leaflet",
         "13 pages: Dashboard, Clinical AI, Imaging, ML Tools, Patients, Predict, Dispatch, Hospitals, Reminders, MCI Board, Ask AI",
         ":5173"),
        (BLUE,   "Backend Tier", "Node.js  ·  Express  ·  Mongoose  ·  Multer  ·  Axios  ·  JWT",
         "/api/data  ·  /api/cluster  ·  /api/ai  ·  /api/dispatch  ·  /api/triage  ·  /api/media  ·  /api/reminders  ·  /api/ml",
         ":5000"),
        (VIOLET, "ML Engine",    "Python 3.11  ·  Flask  ·  scikit-learn  ·  PyTorch  ·  TorchXRayVision  ·  SHAP",
         "/cluster  ·  /risk-profile  ·  /analyze-notes  ·  /forecast-vitals  ·  /analyze-image  ·  /ask  ·  /explain",
         ":5001"),
        (TEAL,   "Storage",      "MongoDB 7  ·  GridFS for images & DICOM",
         "datasets · clusterresults · patients · reminders · dispatches · triages · media",
         ":27017"),
    ]
    for i, (clr, name, stack, routes, port) in enumerate(layers):
        cy = y0 + 0.1 + i * 1.45
        add_rect(s, 0.4, cy, 12.5, 1.25, clr)
        add_text(s, name, 0.6, cy+0.1, 3.3, 0.45, size=16, bold=True, color=WHITE)
        add_text(s, port, 3.9, cy+0.1, 1.2, 0.45, size=14, color=LIGHT, italic=True)
        add_text(s, stack, 0.6, cy+0.55, 12.1, 0.35, size=11, color=RGBColor(0xE0, 0xE7, 0xFF))
        add_text(s, routes, 0.6, cy+0.88, 12.1, 0.35, size=10, color=RGBColor(0xC7, 0xD2, 0xFE), italic=True)
        if i < 3:
            add_text(s, "▼", 6.4, cy+1.22, 1.2, 0.3, size=14, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

def slide_tech_stack():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Technology Stack", "Open-source core, cloud AI optional, no vendor lock-in")
    stack = [
        ("Frontend",   INDIGO, ["React 18", "Vite 5", "TailwindCSS 3", "Recharts + D3.js", "Leaflet.js", "jsPDF", "Axios"]),
        ("Backend",    BLUE,   ["Node.js 25", "Express 4", "Mongoose 8", "Multer", "Anthropic SDK", "JWT", "CORS"]),
        ("ML / AI",    VIOLET, ["Python 3.11", "Flask 3", "scikit-learn 1.4", "SciPy 1.12", "PyTorch 2.2", "TorchXRayVision 1.0", "ONNX"]),
        ("NLP",        TEAL,   ["spaCy 3.7", "scispaCy 0.5", "HuggingFace 4.40", "BioBERT", "Regex NER", "Spell correction"]),
        ("Data + AI",  AMBER,  ["MongoDB 7", "GridFS", "SHAP 0.45", "UMAP-learn", "Prophet 1.1", "FAISS 1.8", "pydicom"]),
    ]
    for i, (title, clr, items) in enumerate(stack):
        cx = 0.4 + i * 2.55
        card(s, cx, y0+0.1, 2.4, 5.4, title, items, color=clr, body_size=11)

def slide_data_flow():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Data Flow", "End-to-end path from CSV upload to clinician insight")
    steps = [
        ("Upload",    "CSV / JSON / sample dataset", INDIGO),
        ("Validate",  "Numeric columns, missing values, IQR outliers", BLUE),
        ("Preprocess","Median impute → StandardScaler → PCA (2D)", VIOLET),
        ("Cluster",   "K-Means · DBSCAN · Hierarchical · GMM", SKY),
        ("Score",     "Composite clinical risk 0-100", AMBER),
        ("Explain",   "SHAP, anomalies, disease risk radar", RED),
        ("Store",     "MongoDB: clusterresults, patients, media", TEAL),
        ("Visualize", "Dashboard, scatter, table, profile, PDF", GREEN),
    ]
    box_w = 1.4
    box_h = 1.5
    gap = 0.1
    total_w = len(steps) * box_w + (len(steps)-1) * gap
    start_x = (13.333 - total_w) / 2
    for i, (name, desc, clr) in enumerate(steps):
        cx = start_x + i * (box_w + gap)
        add_round(s, cx, y0+1.5, box_w, box_h, clr)
        add_text(s, str(i+1), cx, y0+1.55, box_w, 0.3, size=11, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
        add_text(s, name, cx, y0+1.85, box_w, 0.4, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, desc, cx+0.05, y0+2.3, box_w-0.1, 0.7, size=8, color=LIGHT, align=PP_ALIGN.CENTER, italic=True)
        if i < len(steps)-1:
            ax = cx + box_w + 0.005
            add_text(s, "▶", ax, y0+2.0, gap, 0.4, size=14, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    # outcomes row
    out_y = y0 + 3.5
    add_rect(s, 0.4, out_y, 12.5, 1.7, WHITE, line=BORDER)
    add_text(s, "OUTPUTS", 0.6, out_y+0.1, 4, 0.35, size=12, bold=True, color=INDIGO)
    outputs = [
        ("Risk Tier",       "Low / Moderate / High / Critical", RISK_LOW),
        ("Cluster Profile", "Centroid features + size",          BLUE),
        ("SHAP Drivers",    "Top-3 features per patient",        VIOLET),
        ("Anomaly Flags",   "Isolation Forest outliers",        AMBER),
        ("Care Plan",       "Disease-specific actions",          TEAL),
    ]
    for i, (n, d, c) in enumerate(outputs):
        cx = 0.55 + i * 2.5
        add_text(s, n, cx, out_y+0.55, 2.4, 0.32, size=12, bold=True, color=c)
        add_text(s, d, cx, out_y+0.92, 2.4, 0.7, size=10, color=SLATE)
    # key insight
    add_text(s, "Async anomaly detection runs in background after clustering — primary response never blocks on it.",
             0.4, y0+5.5, 12.5, 0.5, size=11, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SECTION 3 — ML ENGINE
# ═════════════════════════════════════════════════════════════════════════════
def slide_preprocessing():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Preprocessing Pipeline", "Every patient row passes through 7 deterministic steps")
    steps = [
        ("1", "Patient ID extraction", "Separate patient_id and sequential id columns from features"),
        ("2", "Numeric column select", "Keep only numeric dtypes; record dropped non-numeric in warnings"),
        ("3", "Median imputation",     "Fill missing values with column median — robust to clinical outliers"),
        ("4", "IQR outlier removal",   "Drop rows outside [Q1 - 1.5·IQR, Q3 + 1.5·IQR] per feature"),
        ("5", "Standard scaling",      "StandardScaler: zero mean, unit variance for distance-based algos"),
        ("6", "PCA reduction",         "Reduce to 2 components for scatter visualization only"),
        ("7", "Report",                "Return rows_before, rows_after, dropped_rows, feature_names"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        col = i % 4
        row = i // 4
        cx = 0.4 + col * 3.2
        cy = y0 + 0.2 + row * 2.5
        add_rect(s, cx, cy, 3.0, 2.2, WHITE, line=BORDER)
        add_rect(s, cx, cy, 3.0, 0.6, INDIGO)
        add_text(s, num, cx+0.12, cy+0.1, 0.5, 0.4, size=22, bold=True, color=WHITE)
        add_text(s, title, cx+0.65, cy+0.15, 2.3, 0.4, size=12, bold=True, color=WHITE)
        add_text(s, desc, cx+0.15, cy+0.75, 2.75, 1.35, size=11, color=SLATE)
    # why callout
    add_rect(s, 0.4, 5.6, 12.5, 1.4, INDIGO)
    add_text(s, "WHY MEDIAN, NOT MEAN?", 0.6, 5.7, 4, 0.3, size=11, bold=True, color=WHITE)
    add_text(s, "Clinical data has extreme outliers (e.g. glucose 9999). A single erroneous value would distort mean imputation for the entire cohort. Median is robust by definition. IQR outlier removal flags suspicious rows for clinical review rather than silently keeping them.",
             0.6, 6.0, 12.1, 1.0, size=12, color=LIGHT)

def slide_clustering_algos():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Four Clustering Algorithms", "From /cluster endpoint — pick one or run all four in parallel")
    algos = [
        ("K-Means", INDIGO, [
            "k-means++ init, 300 max_iter, 10 n_init",
            "Spherical clusters of similar size",
            "Use /optimal-k first for K selection",
            "Output: labels, centroids, inertia",
        ]),
        ("DBSCAN", BLUE, [
            "eps neighborhood, min_samples density",
            "Arbitrary shape clusters",
            "Auto-detects noise as cluster -1",
            "Noise patients get risk_tier = 'Noise'",
        ]),
        ("Hierarchical", VIOLET, [
            "Ward / complete / average / single linkage",
            "Builds dendrogram from linkage_matrix",
            "Post-hoc K via cut-height slider in UI",
            "Best for small cohorts and exploration",
        ]),
        ("GMM", SKY, [
            "Soft assignment via posterior probability",
            "Full / tied / diag / spherical covariance",
            "Returns NxK probability matrix per patient",
            "Best for overlapping elliptical clusters",
        ]),
    ]
    for i, (name, clr, pts) in enumerate(algos):
        cx = 0.4 + i * 3.2
        card(s, cx, y0+0.1, 3.0, 3.3, name, pts, color=clr, body_size=11)
    # comparison + metrics
    add_rect(s, 0.4, y0+3.5, 12.5, 2.4, WHITE, line=BORDER)
    add_text(s, "Algorithm Comparison Mode (algorithm: 'all')", 0.6, y0+3.6, 12, 0.35, size=14, bold=True, color=INDIGO)
    add_text(s, "Run all 4 in parallel, return side-by-side metrics for radar chart in the Comparison Panel.",
             0.6, y0+3.95, 12, 0.3, size=11, italic=True, color=MUTED)
    # metrics
    metrics_y = y0+4.3
    metrics = [
        ("Silhouette", "-1 to 1, higher = better", "Compactness vs separation"),
        ("Davies-Bouldin", "Lower = better", "Avg similarity of cluster to its nearest neighbor"),
        ("Calinski-Harabasz", "Higher = better", "Variance ratio (between/within)"),
    ]
    for i, (n, rng, desc) in enumerate(metrics):
        cx = 0.6 + i * 4.1
        add_round(s, cx, metrics_y, 3.9, 1.5, BG, line=BORDER)
        add_text(s, n, cx+0.15, metrics_y+0.1, 3.7, 0.32, size=12, bold=True, color=VIOLET)
        add_text(s, rng, cx+0.15, metrics_y+0.45, 3.7, 0.3, size=10, color=AMBER, italic=True)
        add_text(s, desc, cx+0.15, metrics_y+0.78, 3.7, 0.65, size=10, color=SLATE)

def slide_risk_scoring():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Risk Scoring Engine", "4-tier clinical risk from composite score, no labels required")
    # tier table
    add_text(s, "RISK TIER THRESHOLDS", 0.4, y0, 12.5, 0.35, size=14, bold=True, color=INDIGO)
    headers = ["Score", "Tier", "Triage", "Action"]
    rows = [
        ["0-24",  "Low",      "P4 - Routine",   "Standard monitoring"],
        ["25-49", "Moderate", "P3 - Soon",      "Follow-up review"],
        ["50-74", "High",     "P2 - Urgent",    "Same-day doctor review"],
        ["75-100","Critical", "P1 - Immediate", "Emergency response"],
    ]
    table(s, 0.4, y0+0.45, 12.5, 2.2, headers, rows,
          col_widths=[1.0, 1.4, 1.8, 4.0], font_size=12, header_size=12)
    # disease-specific
    add_text(s, "SIX DISEASE-SPECIFIC RISK SCORES (computed simultaneously)", 0.4, y0+2.85, 12.5, 0.35, size=13, bold=True, color=VIOLET)
    diseases = [
        ("Diabetes",    "HbA1c ≥ 6.5 (+18-30), glucose > 200 (+12-20), BMI ≥ 30 (+8)"),
        ("Cardiac",     "Troponin > 0.04 (+35), chest pain (+20), hypertensive BP (+15)"),
        ("Kidney",      "Creatinine ≥ 1.5 (+14-24), eGFR < 60 (+16-30), comorbidity (+8)"),
        ("Respiratory", "SpO2 < 94 (+20-35), RR ≥ 24 (+18), COPD history (+18)"),
        ("Stroke",      "Age ≥ 65 (+8), focal neuro symptoms (+18), HTN (+8), AFib (+16)"),
        ("Sepsis",      "WBC abnormal (+16), lactate ≥ 2 (+16-28), SBP < 100 (+14)"),
    ]
    for i, (n, desc) in enumerate(diseases):
        col = i % 3
        row = i // 3
        cx = 0.4 + col * 4.3
        cy = y0+3.3 + row*1.55
        add_round(s, cx, cy, 4.1, 1.4, WHITE, line=BORDER)
        add_rect(s, cx, cy, 0.12, 1.4, VIOLET)
        add_text(s, n, cx+0.25, cy+0.12, 3.7, 0.35, size=12, bold=True, color=VIOLET)
        add_text(s, desc, cx+0.25, cy+0.5, 3.75, 0.85, size=10, color=SLATE)

def slide_shap():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "SHAP Explainability", "Why was this patient assigned to that cluster?")
    add_text(s, "Clustering is unsupervised — no target variable for SHAP. MediCluster uses a clever proxy approach:",
             0.4, y0, 12.5, 0.4, size=12, color=MUTED, italic=True)
    # 4-step flow
    flow = [
        ("1. Proxy model",     "Train Ridge regression to predict cluster ID from features"),
        ("2. LinearExplainer", "SHAP applied with feature_perturbation='interventional'"),
        ("3. Global importance", "Mean |SHAP| per feature across all patients"),
        ("4. Per-patient",     "Top-3 drivers with direction (↑ or ↓ risk)"),
    ]
    for i, (title, desc) in enumerate(flow):
        cx = 0.4 + i * 3.2
        add_round(s, cx, y0+0.7, 3.0, 1.7, INDIGO)
        add_text(s, title, cx, y0+0.85, 3.0, 0.4, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, desc, cx+0.15, y0+1.3, 2.7, 1.05, size=10, color=LIGHT, align=PP_ALIGN.CENTER)
    # interpreting callout
    add_rect(s, 0.4, y0+2.7, 12.5, 1.5, WHITE, line=BORDER)
    add_text(s, "INTERPRETING SHAP VALUES", 0.6, y0+2.8, 4, 0.35, size=13, bold=True, color=VIOLET)
    bullets = [
        "Positive SHAP value → feature pushes patient toward higher-numbered (higher-risk) clusters",
        "Negative SHAP value → feature pulls patient toward lower-risk clusters",
        "High mean |SHAP| → feature strongly separates clusters globally",
        "Subsampled to ≤200 patients for performance; configurable via max_samples",
    ]
    bullet_block(s, bullets, 0.6, y0+3.15, 12.1, 1.0, size=11, color=SLATE)
    # graceful degradation
    add_rect(s, 0.4, y0+4.3, 12.5, 1.3, AMBER)
    add_text(s, "GRACEFUL DEGRADATION", 0.6, y0+4.4, 4, 0.3, size=11, bold=True, color=WHITE)
    add_text(s, "If SHAP is not installed, endpoint returns zero scores with a stub message. Frontend hides the SHAP panel automatically. The system never fails — it always returns something useful.",
             0.6, y0+4.7, 12.1, 0.9, size=11, color=WHITE)

def slide_anomalies():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Anomaly Detection", "Isolation Forest flags patients who don't fit the cohort")
    add_text(s, "Statistically unusual patients may indicate rare disease, data entry errors, or candidates for individual clinical review. Auto-triggered after every clustering run, async, non-blocking.",
             0.4, y0, 12.5, 0.6, size=12, color=MUTED, italic=True)
    # algorithm box
    add_round(s, 0.4, y0+0.85, 6.0, 3.3, WHITE, line=BORDER)
    add_rect(s, 0.4, y0+0.85, 6.0, 0.5, VIOLET)
    add_text(s, "ISOLATION FOREST ALGORITHM", 0.55, y0+0.92, 5.7, 0.35, size=13, bold=True, color=WHITE)
    algo = [
        "Build random decision trees by selecting random feature & split value",
        "Recursively partition data until each point is isolated",
        "Fewer splits to isolate → more anomalous (shorter path length)",
        "Anomaly score = -2^(-avg_path_length / c(n))",
        "contamination = 0.05 (default 5% expected anomalies)",
        "Returns: anomaly_flags (-1 / 1), scores, count, contributions",
    ]
    bullet_block(s, algo, 0.55, y0+1.45, 5.7, 2.7, size=11, color=SLATE)
    # frontend
    add_round(s, 6.7, y0+0.85, 6.2, 3.3, WHITE, line=BORDER)
    add_rect(s, 6.7, y0+0.85, 6.2, 0.5, RED)
    add_text(s, "FRONTEND: ANOMALY WATCHLIST", 6.85, y0+0.92, 5.9, 0.35, size=13, bold=True, color=WHITE)
    ui = [
        "Red-bordered list in Dashboard 'Anomalies' tab",
        "Shows anomaly score, risk tier, top feature values",
        "Click any patient → Patient Detail Modal opens",
        "Empty state: green checkmark with 'No anomalies'",
        "Async run: results merge into dashboard when ready",
        "No timeout risk on primary clustering response",
    ]
    bullet_block(s, ui, 6.85, y0+1.45, 5.9, 2.7, size=11, color=SLATE)
    # when to use
    add_rect(s, 0.4, y0+4.4, 12.5, 1.4, INDIGO)
    add_text(s, "WHEN TO INVESTIGATE A FLAGGED PATIENT", 0.6, y0+4.5, 6, 0.3, size=11, bold=True, color=WHITE)
    add_text(s, "Rare disease presentation  ·  Data entry error to verify  ·  Patient who doesn't fit any standard cluster  ·  Outlier feature values (e.g. troponin 50x normal)  ·  Reasonable clinical scenario the model hasn't seen before",
             0.6, y0+4.8, 12.1, 0.9, size=11, color=LIGHT)

def slide_automl():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "AutoML & Feature Selection", "Smarter K, smarter features, smarter predictions")
    features = [
        ("Optimal K Finder", INDIGO, [
            "Elbow + Silhouette sweep K=2..10",
            "Returns recommended K with rationale",
            "Use silhouette K if score > 0.3, else elbow",
        ]),
        ("Feature Importance", VIOLET, [
            "Mutual Information (non-linear)",
            "ANOVA F-Score (linear separation)",
            "Composite: 0.6·MI + 0.4·F (clinical non-linearity weight)",
        ]),
        ("UMAP / t-SNE", BLUE, [
            "Better 2D structure than PCA",
            "UMAP: n_neighbors=15, min_dist=0.1",
            "t-SNE: perplexity auto-clamped to (N-1)/3",
        ]),
        ("Supervised AutoML", TEAL, [
            "Auto-detect classification vs regression",
            "3 models: LogReg / RandomForest / GradientBoosting",
            "F1 leaderboard for classification, R² for regression",
        ]),
    ]
    for i, (title, clr, pts) in enumerate(features):
        cx = 0.4 + i * 3.2
        card(s, cx, y0+0.1, 3.0, 3.0, title, pts, color=clr, body_size=11)
    # targets
    add_rect(s, 0.4, y0+3.3, 12.5, 2.4, WHITE, line=BORDER)
    add_text(s, "COMMON SUPERVISED TARGETS", 0.6, y0+3.4, 6, 0.35, size=13, bold=True, color=INDIGO)
    targets = [
        ("risk_label",       "Predicted risk tier (classification)"),
        ("readmitted",       "30-day readmission (classification)"),
        ("icu_admission",    "ICU need prediction (classification)"),
        ("length_of_stay",   "Predicted LOS in days (regression)"),
        ("mortality_risk",   "In-hospital mortality (classification)"),
    ]
    for i, (col_name, desc) in enumerate(targets):
        col = i % 3
        row = i // 3
        cx = 0.6 + col * 4.1
        cy = y0+3.8 + row*0.85
        add_round(s, cx, cy, 3.9, 0.75, BG, line=BORDER)
        add_text(s, col_name, cx+0.15, cy+0.1, 3.7, 0.3, size=11, bold=True, color=AMBER, font="Consolas")
        add_text(s, desc, cx+0.15, cy+0.42, 3.7, 0.3, size=10, color=SLATE)

# ═════════════════════════════════════════════════════════════════════════════
# SECTION 4 — CLINICAL AI
# ═════════════════════════════════════════════════════════════════════════════
def slide_nlp():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Clinical Notes NLP", "20+ structured fields extracted from free-text clinical notes")
    # pipeline as horizontal steps
    pipe = [
        "Spell\ncorrection",
        "Risk keywords\n(4-tier)",
        "Trajectory\ndetection",
        "Drug\nextraction",
        "Lab value\nextraction",
        "Symptom\nextraction",
        "SpaCy\nNER",
        "ICD-10\nsuggest",
        "Risk tier\ninfer",
        "Department\nrouting",
    ]
    box_w = 1.18
    gap = 0.07
    start_x = (13.333 - (len(pipe)*box_w + (len(pipe)-1)*gap)) / 2
    for i, step in enumerate(pipe):
        cx = start_x + i * (box_w + gap)
        add_round(s, cx, y0+0.1, box_w, 0.95, INDIGO)
        add_text(s, step, cx+0.05, y0+0.2, box_w-0.1, 0.8, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # details
    details_y = y0+1.25
    cards = [
        ("Symptom Extraction", VIOLET, [
            "21 canonical symptoms tracked",
            "Negation detection (45-char window)",
            "Severity & duration extraction",
            "Positive mention wins over negated",
        ]),
        ("ICD-10 Mapping", INDIGO, [
            "Rule-based, 19 common conditions",
            "Pneumonia, sepsis, MI, stroke, AKI...",
            "Keyword → code mapping",
            "Returns code + description",
        ]),
        ("Risk Keyword Lexicon", RED, [
            "Critical: cardiac arrest, ICU, intubated",
            "High: PE, pneumonia, AKI, deteriorating",
            "Moderate: diabetes, HTN, AFib",
            "Low: stable, improving, discharged",
        ]),
        ("Trajectory + Summary", TEAL, [
            "Improving / stable / deteriorating",
            "Compares improvement vs deterioration hits",
            "DistilBART abstractive summary (optional)",
            "Template fallback when no transformer",
        ]),
    ]
    for i, (t, c, pts) in enumerate(cards):
        cx = 0.4 + i * 3.2
        card(s, cx, details_y, 3.0, 2.5, t, pts, color=c, body_size=10)
    # callout
    add_rect(s, 0.4, y0+4.7, 12.5, 1.0, BG, line=BORDER)
    add_text(s, "NLP BACKEND DETECTION", 0.6, y0+4.8, 4, 0.3, size=11, bold=True, color=INDIGO)
    add_text(s, "Response includes nlp_backend: 'scispacy' or 'regex-fallback' so the frontend can show sophistication tier to the clinician.",
             0.6, y0+5.1, 12.1, 0.6, size=11, color=SLATE)

def slide_forecasting():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Vitals Forecasting + MEWS", "Predict deterioration 3+ readings ahead, validated clinical early-warning score")
    # MEWS table
    add_text(s, "MEWS (MODIFIED EARLY WARNING SCORE)", 0.4, y0, 6, 0.35, size=14, bold=True, color=INDIGO)
    mews_headers = ["MEWS", "Alert", "Action"]
    mews_rows = [
        ["0-1",  "Low",      "Routine monitoring every 12 hours"],
        ["2-3",  "Moderate", "Notify charge nurse, monitor every 4 hours"],
        ["4-5",  "High",     "Urgent medical review within 30 minutes"],
        ["≥6",   "Critical", "IMMEDIATE medical emergency response"],
    ]
    table(s, 0.4, y0+0.45, 6.0, 2.2, mews_headers, mews_rows, col_widths=[1.0, 1.0, 4.0], font_size=11, header_size=11)
    # methods
    add_text(s, "FORECASTING METHODS", 6.7, y0, 6.2, 0.35, size=14, bold=True, color=VIOLET)
    methods = [
        ("LSTM",      "PyTorch 2-layer LSTM, hidden=32, 60 epochs Adam lr=0.01"),
        ("Prophet",   "Trend + seasonality decomposition, ≥10 readings"),
        ("Linear",    "numpy.polyfit fallback, ±0.5·std confidence band"),
        ("Auto",      "Prophet if available & ≥10 readings, else LSTM, else Linear"),
    ]
    for i, (name, desc) in enumerate(methods):
        cy = y0+0.5 + i*0.55
        add_round(s, 6.7, cy, 6.2, 0.5, BG, line=BORDER)
        add_rect(s, 6.7, cy, 1.0, 0.5, VIOLET)
        add_text(s, name, 6.7, cy+0.1, 1.0, 0.3, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, desc, 7.8, cy+0.1, 5.0, 0.3, size=10, color=SLATE)
    # deterioration risk rules
    add_text(s, "DETERIORATION RISK ASSESSMENT", 0.4, y0+2.95, 12.5, 0.35, size=14, bold=True, color=RED)
    add_text(s, "After forecasting, trend-based signals are counted and combined with MEWS:",
             0.4, y0+3.3, 12.5, 0.3, size=11, italic=True, color=MUTED)
    rules = [
        "HR or RR forecast > 10% above last value → +1 deterioration signal",
        "Systolic BP forecast > 10% below last value → +1 deterioration signal",
        "MEWS ≥ 6 OR ≥ 3 trend signals → 'critical'",
        "MEWS ≥ 4 OR ≥ 2 trend signals → 'high'",
        "MEWS ≥ 2 OR ≥ 1 signal → 'moderate'",
        "Otherwise → 'low'",
    ]
    bullet_block(s, rules, 0.4, y0+3.65, 12.5, 2.0, size=12, color=SLATE)
    add_rect(s, 0.4, y0+5.7, 12.5, 0.9, INDIGO)
    add_text(s, "MINIMUM REQUIREMENTS", 0.6, y0+5.8, 4, 0.3, size=11, bold=True, color=WHITE)
    add_text(s, "4 data points for LSTM  ·  10 for Prophet  ·  Forecast returns yhat, yhat_lower, yhat_upper confidence intervals",
             0.6, y0+6.1, 12.1, 0.5, size=11, color=LIGHT)

def slide_imaging():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Medical Imaging AI", "Chest X-ray pathology detection with Grad-CAM explainability")
    # 5 models
    add_text(s, "FIVE TORCHXRAYVISION MODELS", 0.4, y0, 12.5, 0.35, size=14, bold=True, color=INDIGO)
    models = [
        ("densenet121-res224-chex",  "DenseNet-121", "Stanford CheXpert (224k)",   "224"),
        ("densenet121-res224-all",   "DenseNet-121", "CheXpert+NIH+MIMIC+PadChest","224"),
        ("resnet50-res512-all",      "ResNet-50",    "All datasets",               "512"),
        ("densenet121-res224-nih",   "DenseNet-121", "NIH ChestX-ray14",           "224"),
        ("densenet121-res224-pc",    "DenseNet-121", "PadChest",                   "224"),
    ]
    headers = ["Model ID", "Architecture", "Training Data", "Input"]
    rows = [list(m) for m in models]
    table(s, 0.4, y0+0.4, 12.5, 2.0, headers, rows, col_widths=[3.5, 2.0, 5.0, 1.5], font_size=11, header_size=11)
    add_text(s, "Recommended default: densenet121-res224-chex (CheXNet) — best calibrated for 14-class detection",
             0.4, y0+2.45, 12.5, 0.3, size=10, italic=True, color=MUTED)
    # pipeline
    add_text(s, "PREPROCESSING → INFERENCE → CLINICAL ENRICHMENT", 0.4, y0+2.85, 12.5, 0.35, size=14, bold=True, color=VIOLET)
    add_round(s, 0.4, y0+3.3, 12.5, 1.7, WHITE, line=BORDER)
    pipe = [
        "DICOM detection\n(pydicom)",
        "CLAHE enhancement\n(non-DICOM only)",
        "TorchXRayVision\nnormalize [-1024, 1024]",
        "Center crop\n+ resize",
        "Forward pass\n+ sigmoid",
        "Filter ≥15%\nconfidence",
        "Attach clinical\ncontext",
    ]
    box_w = (12.5 - (len(pipe)-1)*0.05) / len(pipe)
    for i, step in enumerate(pipe):
        cx = 0.4 + i * (box_w + 0.05)
        add_rect(s, cx, y0+3.45, box_w, 1.4, BG, line=BORDER)
        add_text(s, step, cx+0.05, y0+3.55, box_w-0.1, 1.2, size=10, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
    # gradcam
    add_rect(s, 0.4, y0+5.2, 12.5, 1.4, RED)
    add_text(s, "GRAD-CAM HEATMAPS", 0.6, y0+5.3, 4, 0.3, size=11, bold=True, color=WHITE)
    add_text(s, "Generate visual evidence of where the model focused for each pathology. Red regions = strongest activation. Returns base64 PNG for overlay on the original image. Downloadable as PNG.",
             0.6, y0+5.6, 12.1, 0.9, size=11, color=WHITE)

def slide_rag_drug():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "RAG Chatbot + Drug Interaction Checker", "Natural language Q&A over indexed patients, with built-in medication safety")
    # RAG left
    add_round(s, 0.4, y0+0.1, 6.2, 5.5, WHITE, line=BORDER)
    add_rect(s, 0.4, y0+0.1, 6.2, 0.5, VIOLET)
    add_text(s, "RAG CHATBOT — ASK AI", 0.55, y0+0.17, 5.9, 0.35, size=13, bold=True, color=WHITE)
    rag = [
        "User query: 'Which High Risk patients are over 60?'",
        "TF-IDF vector search (top-50 matching patients)",
        "Structured filter extraction (risk_tier, age range, cluster_id)",
        "Filter application to retrieved patients",
        "Answer generation: stats + sample patient list",
        "Returns {answer, matched_patients, filters_applied}",
        "FAISS optional: replace exact cosine with ANN",
        "Index rebuilt after every /cluster run",
    ]
    bullet_block(s, rag, 0.55, y0+0.7, 5.9, 4.8, size=11, color=SLATE)
    # Drug right
    add_round(s, 6.8, y0+0.1, 6.1, 5.5, WHITE, line=BORDER)
    add_rect(s, 6.8, y0+0.1, 6.1, 0.5, RED)
    add_text(s, "DRUG INTERACTION CHECKER", 6.95, y0+0.17, 5.8, 0.35, size=13, bold=True, color=WHITE)
    drug_headers = ["Pair", "Severity"]
    drug_rows = [
        ["Warfarin + Aspirin",          "Major"],
        ["Warfarin + Ibuprofen",        "Major"],
        ["Metformin + Alcohol",         "Major"],
        ["Metoprolol + Verapamil",      "Major"],
        ["SSRI + Tramadol",             "Major"],
        ["Lithium + Ibuprofen",         "Major"],
        ["Lisinopril + Potassium",      "Moderate"],
        ["Furosemide + Digoxin",        "Moderate"],
    ]
    table(s, 6.95, y0+0.7, 5.8, 3.6, drug_headers, drug_rows, col_widths=[3.5, 1.5], font_size=10, header_size=10)
    add_text(s, "12 high-priority pairs in knowledge base. Substring matching handles brand names, abbreviations, partial matches. Auto-triggers when 2+ meds extracted from notes.",
             6.95, y0+4.4, 5.8, 1.2, size=10, italic=True, color=MUTED)

# ═════════════════════════════════════════════════════════════════════════════
# SECTION 5 — FRONTEND PAGES
# ═════════════════════════════════════════════════════════════════════════════
def slide_pages():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Frontend — 13 Pages", "Single-page React app, top-nav layout, dark clinical aesthetic")
    pages = [
        ("Dashboard",     "Upload, configure, scatter, risk donut, metrics, SHAP, PDF export",   INDIGO),
        ("Clinical AI",   "Real-time NLP, risk profile, MEWS, care plan, drug check, vitals",   VIOLET),
        ("Imaging",       "X-ray upload, 5 models, findings, Grad-CAM, clinical context",      BLUE),
        ("ML Tools",      "Optimal K, UMAP/t-SNE, supervised AutoML, feature importance",       TEAL),
        ("Patients",      "Search, media, cluster history, similar patient search",            AMBER),
        ("Predict",       "Single-patient risk from saved cluster centroids",                   SKY),
        ("Dispatch",      "Triage assessment, ambulance dispatch, hospital routing",           RED),
        ("Hospitals",     "Leaflet map, Overpass query, OSRM routing, distance & capacity",     GREEN),
        ("Reminders",     "Clinical task scheduler, priority, recurring, PDF/CSV export",       INDIGO),
        ("MCI Board",     "Mass Casualty grid by START triage (black/red/yellow/green)",        VIOLET),
        ("Ask AI",        "Chat interface, complex patient queries, applied filters",          BLUE),
        ("Home / Login",  "Landing, onboarding, 4 algorithm hero cards, 3-step demo",           TEAL),
    ]
    # 4 rows × 3 cols, row 4 has 1 tile (Home/Login) on left, rest empty
    for i, (page, desc, clr) in enumerate(pages):
        col = i % 3
        row = i // 3
        cx = 0.4 + col * 4.3
        cy = y0 + 0.15 + row * 1.4
        add_round(s, cx, cy, 4.1, 1.3, WHITE, line=BORDER)
        add_rect(s, cx, cy, 0.12, 1.3, clr)
        add_text(s, page, cx+0.25, cy+0.1, 3.8, 0.35, size=13, bold=True, color=clr)
        add_text(s, desc, cx+0.25, cy+0.5, 3.8, 0.75, size=10, color=SLATE)
    # legend chip below
    add_rect(s, 0.4, y0+5.85, 12.5, 0.35, BG, line=BORDER)
    add_text(s, "Design: Deep navy (#0A0F1E) + electric teal (#00D4AA). Monospace for vitals, sans-serif for UI. Flat dark panels, no gradients.",
             0.6, y0+5.88, 12.1, 0.3, size=9, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

def slide_dispatch():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Ambulance Dispatch + MCI Triage", "Real-time emergency coordination with mass casualty management")
    # dispatch left
    add_round(s, 0.4, y0+0.1, 6.2, 5.4, WHITE, line=BORDER)
    add_rect(s, 0.4, y0+0.1, 6.2, 0.5, INDIGO)
    add_text(s, "ARIA — AMBULANCE DISPATCH", 0.55, y0+0.17, 5.9, 0.35, size=13, bold=True, color=WHITE)
    dispatch_items = [
        "Live Leaflet.js map with ambulance positions",
        "Real-time status: available / en_route / at_scene / transporting",
        "WebSocket events: dispatch:new, driver:status, location:update",
        "OSRM road-network routing (replaces Haversine straight-line)",
        "Claude selects nearest ambulance + best hospital by bed count and specialty",
        "Route polyline visualization on map",
        "Driver dashboard with patient alert overlays",
        "Two-way dispatcher-driver communication",
        "Full dispatch history (resolve/cancel)",
    ]
    bullet_block(s, dispatch_items, 0.55, y0+0.7, 5.9, 4.7, size=11, color=SLATE)
    # MCI right
    add_round(s, 6.8, y0+0.1, 6.1, 5.4, WHITE, line=BORDER)
    add_rect(s, 6.8, y0+0.1, 6.1, 0.5, RED)
    add_text(s, "MCI TRIAGE BOARD", 6.95, y0+0.17, 5.8, 0.35, size=13, bold=True, color=WHITE)
    # colors row
    colors = [
        ("Black",  "Deceased/Expectant",  RGBColor(0x1F, 0x29, 0x37)),
        ("Red",    "Immediate",           RED),
        ("Yellow", "Delayed",             AMBER),
        ("Green",  "Minor",               GREEN),
    ]
    for i, (color, desc, c) in enumerate(colors):
        cx = 6.95 + (i % 2) * 2.95
        cy = y0+0.75 + (i // 2) * 0.7
        add_rect(s, cx, cy, 2.8, 0.6, c)
        add_text(s, color, cx, cy+0.1, 2.8, 0.25, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, desc, cx, cy+0.32, 2.8, 0.25, size=9, color=LIGHT, align=PP_ALIGN.CENTER)
    mci = [
        "Voice-to-triage: speak/type patient condition",
        "Claude extracts vitals + assigns P1/P2/P3/P4 priority",
        "Stateless scoring: instant score, no DB write",
        "Hospital allocation by triage category",
        "Pathway flags (trauma / cardiac / respiratory)",
        "Live patient cards with vitals history",
        "Transport status tracking per patient",
    ]
    bullet_block(s, mci, 6.95, y0+2.25, 5.85, 3.2, size=11, color=SLATE)

# ═════════════════════════════════════════════════════════════════════════════
# SECTION 6 — API REFERENCE
# ═════════════════════════════════════════════════════════════════════════════
def slide_api():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "API Reference (25+ endpoints)", "ML engine on :5001, backend on :5000")
    # table of endpoints
    headers = ["Endpoint", "Purpose", "Method"]
    rows = [
        ["/cluster",            "Run clustering (kmeans/dbscan/hierarchical/gmm/all)",        "POST"],
        ["/preprocess-preview", "Dataset stats + dropped-row warnings",                          "POST"],
        ["/optimal-k",          "Elbow + Silhouette sweep, recommended K",                       "POST"],
        ["/explain",            "SHAP global + per-patient feature importance",                 "POST"],
        ["/detect-anomalies",   "Isolation Forest outlier flags",                                "POST"],
        ["/risk-profile",       "Single-patient composite risk, 6 disease scores, care plan",   "POST"],
        ["/compare-visits",     "Cross-visit risk delta, trend label",                           "POST"],
        ["/similar-patients",   "Top-K clinical similarity search",                             "POST"],
        ["/analyze-notes",      "Full NLP pipeline: trajectory, ICD-10, drugs, symptoms",       "POST"],
        ["/drug-interactions",  "12-pair knowledge base check",                                  "POST"],
        ["/forecast-vitals",    "LSTM/Prophet forecast + MEWS + deterioration",                "POST"],
        ["/analyze-image",      "Chest X-ray pathology detection (5 models)",                    "POST"],
        ["/gradcam",            "Heatmap overlay for image model focus",                        "POST"],
        ["/ask",                "RAG chatbot natural language Q&A",                              "POST"],
    ]
    table(s, 0.4, y0+0.1, 12.5, 5.5, headers, rows, col_widths=[2.5, 8.0, 1.0], font_size=10, header_size=11)
    add_text(s, "Backend proxies all ML calls with 120s timeout; persists results to MongoDB for history comparison.",
             0.4, y0+5.65, 12.5, 0.4, size=10, italic=True, color=MUTED)

# ═════════════════════════════════════════════════════════════════════════════
# SECTION 7 — DEPLOYMENT, PERFORMANCE, SECURITY, TESTING
# ═════════════════════════════════════════════════════════════════════════════
def slide_deployment():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Deployment & Quick Start", "Docker Compose for one-command bring-up")
    # code block
    add_rect(s, 0.4, y0+0.1, 8.0, 4.0, RGBColor(0x1E, 0x1B, 0x4B))
    add_text(s, "docker-compose.yml (excerpt)", 0.6, y0+0.2, 7.6, 0.3, size=11, bold=True, color=LIGHT)
    code = """services:
  ml-engine:
    build: ./ml-engine
    ports: ["5001:5001"]
    environment:
      - FLASK_ENV=production

  backend:
    build: ./backend
    ports: ["5000:5000"]
    depends_on: [ml-engine, mongodb]
    environment:
      - MONGO_URI=mongodb://mongodb:27017/medicluster
      - ML_ENGINE_URL=http://ml-engine:5001

  frontend:
    build: ./frontend
    ports: ["5173:80"]
    depends_on: [backend]

  mongodb:
    image: mongo:7
    volumes: [mongo_data:/data/db]"""
    add_text(s, code, 0.6, y0+0.55, 7.6, 3.45, size=10, color=LIGHT, font="Consolas")
    # quick start
    add_round(s, 8.6, y0+0.1, 4.3, 4.0, WHITE, line=BORDER)
    add_rect(s, 8.6, y0+0.1, 4.3, 0.45, INDIGO)
    add_text(s, "QUICK START", 8.75, y0+0.17, 4.0, 0.3, size=12, bold=True, color=WHITE)
    qs = [
        "cd medicluster",
        "docker compose up --build",
        "",
        "Or manual:",
        "  ml-engine: python app.py (:5001)",
        "  backend:   npm run dev (:5000)",
        "  frontend:  npm run dev (:5173)",
    ]
    add_text(s, "\n".join(qs), 8.75, y0+0.7, 4.0, 3.3, size=11, color=SLATE, font="Consolas")
    # infra footer
    add_rect(s, 0.4, y0+4.3, 12.5, 1.4, BG, line=BORDER)
    add_text(s, "PRODUCTION INFRASTRUCTURE", 0.6, y0+4.4, 5, 0.3, size=11, bold=True, color=INDIGO)
    infra = [
        "Gunicorn: 2 workers, 300s timeout (LSTM + Grad-CAM)",
        "Ambulance dispatch: Kubernetes manifests + Terraform configs",
        "Nginx proxies /api/ requests from frontend container",
        "Health check endpoints for container orchestration",
    ]
    bullet_block(s, infra, 0.6, y0+4.7, 12.1, 1.0, size=11, color=SLATE)

def slide_performance():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Performance & Scalability", "Measured on commodity hardware, no GPU required")
    headers = ["Operation", "Dataset Size", "Typical Latency"]
    rows = [
        ["K-Means clustering",          "100 patients",          "< 200ms"],
        ["K-Means clustering",          "10,000 patients",       "< 5s"],
        ["NLP notes analysis",          "Short note",            "< 500ms"],
        ["NLP notes analysis",          "Full discharge summary", "< 2s"],
        ["LSTM training + forecast",    "20 readings, 3 steps",  "< 3s"],
        ["Prophet forecast",            "50 readings",           "< 2s"],
        ["Chest X-ray inference",       "Single image",          "< 1s (CPU)"],
        ["SHAP computation",            "200 patients",          "< 10s"],
        ["Optimal K sweep",             "K=2-10, 500 patients",  "< 15s"],
    ]
    table(s, 0.4, y0+0.1, 12.5, 4.0, headers, rows, col_widths=[4.0, 4.0, 4.5], font_size=11, header_size=12)
    # scalability cards
    cards = [
        ("Horizontal scaling", INDIGO, "Stateless ML engine. Multiple instances behind LB; sticky sessions for chatbot."),
        ("Model caching",      VIOLET, "X-ray models cached in memory after first load; subsequent calls are fast."),
        ("SHAP subsampling",   BLUE,   "Computed on ≤200 samples to prevent timeout on large datasets."),
        ("Async anomalies",    TEAL,   "Anomaly detection runs in background, never blocks primary response."),
    ]
    for i, (t, c, desc) in enumerate(cards):
        cx = 0.4 + i * 3.2
        card(s, cx, y0+4.3, 3.0, 1.7, t, [desc], color=c, body_size=10)

def slide_security():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Security & Compliance", "Local processing, no external APIs, audit trail by default")
    add_text(s, "CURRENT", 0.4, y0, 6.0, 0.35, size=14, bold=True, color=INDIGO)
    add_text(s, "RECOMMENDED FOR PRODUCTION", 6.8, y0, 6.0, 0.35, size=14, bold=True, color=RED)
    current = [
        "CORS restricted to localhost:3000, :5173",
        "All API inputs validated before ML processing",
        "50MB max upload size (CSV + images)",
        "Error sanitization in production mode",
        "No secrets hardcoded in source code",
        "Sample data uses synthetic patient records",
        "All ML inference runs locally",
        "Audit trail: all clustering runs timestamped",
    ]
    bullet_block(s, current, 0.4, y0+0.4, 6.0, 5.0, size=11, color=SLATE)
    prod = [
        "HTTPS/TLS termination at load balancer",
        "JWT-based auth for all API routes",
        "Role-based access (physician / nurse / admin)",
        "Field-level encryption for sensitive PHI",
        "HIPAA-compliant audit logging",
        "Data retention policies per local regulation",
        "DICOM metadata anonymization (planned)",
        "EHR connector for HL7 FHIR (Phase 2)",
    ]
    bullet_block(s, prod, 6.8, y0+0.4, 6.0, 5.0, size=11, color=SLATE)
    # disclaimer
    add_rect(s, 0.4, y0+5.5, 12.5, 0.9, AMBER)
    add_text(s, "NOT FDA/CE CLEARED", 0.6, y0+5.6, 3, 0.3, size=11, bold=True, color=WHITE)
    add_text(s, "MediCluster is a decision support tool. All predictions must be reviewed by qualified clinicians. Not a substitute for professional medical judgment.",
             0.6, y0+5.9, 12.1, 0.5, size=11, color=WHITE)

def slide_testing():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Testing Strategy", "Smoke tests, feature tests, integration tests, manual API dashboard")
    cards = [
        ("ML Smoke Tests", INDIGO, [
            "ml-engine/test_smoke.py",
            "Health endpoint returns 'ok'",
            "Clustering response shape correct",
            "Risk profile returns required fields",
        ]),
        ("Feature Tests", VIOLET, [
            "ml-engine/test_new_features.py",
            "NLP with sample clinical text",
            "Forecasting with synthetic time series",
            "Drug detection with known pairs",
            "Anomalies with injected outliers",
        ]),
        ("Integration", TEAL, [
            "ambulance-dispatch/tests/",
            "WebSocket connection establishment",
            "Dispatch event flow (new → complete)",
            "Map rendering with mock coordinates",
        ]),
        ("API Dashboard", AMBER, [
            "Interactive HTML at :5001/pi_tester.html",
            "Full form interface for all 25+ endpoints",
            "No Postman required for manual testing",
            "Live during development",
        ]),
    ]
    for i, (t, c, pts) in enumerate(cards):
        cx = 0.4 + i * 3.2
        card(s, cx, y0+0.1, 3.0, 5.0, t, pts, color=c, body_size=11)

# ═════════════════════════════════════════════════════════════════════════════
# SECTION 8 — ROADMAP, LIMITATIONS, GLOSSARY
# ═════════════════════════════════════════════════════════════════════════════
def slide_roadmap():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Feature Roadmap", "Three phases over the next 18 months")
    # Phase 1: Complete
    add_rect(s, 0.4, y0+0.1, 4.1, 5.6, GREEN)
    add_text(s, "PHASE 1", 0.4, y0+0.25, 4.1, 0.3, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Complete", 0.4, y0+0.55, 4.1, 0.4, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    p1 = [
        "Core clustering (K-Means, DBSCAN, Hierarchical, GMM)",
        "Risk scoring & 4-tier triage",
        "NLP notes analysis (NER, ICD-10, trajectory)",
        "Chest X-ray AI (5 models + Grad-CAM)",
        "Vital sign forecasting (LSTM + Prophet)",
        "RAG chatbot",
        "Drug interaction checker (12 pairs)",
        "SHAP explainability",
        "Anomaly detection (Isolation Forest)",
        "Ambulance dispatch system",
    ]
    bullet_block(s, p1, 0.55, y0+1.1, 3.8, 4.5, size=10, color=WHITE)
    # Phase 2: Planned
    add_rect(s, 4.7, y0+0.1, 4.1, 5.6, AMBER)
    add_text(s, "PHASE 2", 4.7, y0+0.25, 4.1, 0.3, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Planned (6-12 mo)", 4.7, y0+0.55, 4.1, 0.4, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    p2 = [
        "Federated learning across hospital networks",
        "Real-time HL7 FHIR integration",
        "EHR connector (Epic, Cerner, Allscripts)",
        "Mobile app for ward nurses",
        "Voice input for clinical notes",
        "Multilingual clinical notes (Hindi, Tamil, etc.)",
        "Wearable device integration (smartwatch vitals)",
        "Population health trend analysis",
    ]
    bullet_block(s, p2, 4.85, y0+1.1, 3.8, 4.5, size=10, color=WHITE)
    # Phase 3: Future
    add_rect(s, 9.0, y0+0.1, 4.1, 5.6, VIOLET)
    add_text(s, "PHASE 3", 9.0, y0+0.25, 4.1, 0.3, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Future (12-18 mo)", 9.0, y0+0.55, 4.1, 0.4, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    p3 = [
        "Predictive readmission prevention program",
        "Insurance risk stratification module",
        "Genomic risk factor integration",
        "Remote patient monitoring (IoT sensors)",
        "Multi-hospital benchmarking",
        "Continuous model retraining pipeline",
        "Adversarial robustness testing",
        "Multi-modal fusion (text + image + vitals)",
    ]
    bullet_block(s, p3, 9.15, y0+1.1, 3.8, 4.5, size=10, color=WHITE)

def slide_limitations():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Known Limitations", "Honest accounting of where MediCluster is not yet production-grade")
    add_text(s, "CLINICAL", 0.4, y0, 6.0, 0.35, size=14, bold=True, color=RED)
    add_text(s, "TECHNICAL", 6.8, y0, 6.0, 0.35, size=14, bold=True, color=INDIGO)
    clinical = [
        "Not FDA/CE cleared — decision support only, not a medical device",
        "Unsupervised baseline — without labels, risk tiers are heuristic",
        "Limited ICD-10 coverage — ~19 common conditions mapped",
        "Drug interaction coverage — 12 pairs in knowledge base",
        "Requires qualified clinician review of every prediction",
    ]
    bullet_block(s, clinical, 0.4, y0+0.4, 6.0, 5.0, size=12, color=SLATE)
    tech = [
        "In-memory patient index for RAG — resets on server restart",
        "Single-node ML — LSTM trains per-request, may degrade under high concurrency",
        "DICOM metadata — full anonymization & multi-frame not implemented",
        "No model persistence — supervised AutoML trained per request",
        "FAISS optional — falls back to exact cosine which is slow for large cohorts",
    ]
    bullet_block(s, tech, 6.8, y0+0.4, 6.0, 5.0, size=12, color=SLATE)
    add_rect(s, 0.4, y0+5.5, 12.5, 0.9, VIOLET)
    add_text(s, "WHAT WE MITIGATE", 0.6, y0+5.6, 3, 0.3, size=11, bold=True, color=WHITE)
    add_text(s, "Graceful degradation at every layer: SHAP → zero stubs, LSTM → linear extrapolation, SpaCy → regex NER, FAISS → exact search. The system never fails — it always returns useful output.",
             0.6, y0+5.9, 12.1, 0.5, size=11, color=LIGHT)

def slide_glossary():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Glossary", "Key clinical and technical terms used throughout the deck")
    terms = [
        ("AVPU",          "Alert, Voice, Pain, Unresponsive — consciousness scale"),
        ("BUN",           "Blood Urea Nitrogen"),
        ("CLAHE",         "Contrast Limited Adaptive Histogram Equalization"),
        ("DBSCAN",        "Density-Based Spatial Clustering of Applications with Noise"),
        ("eGFR",          "Estimated Glomerular Filtration Rate"),
        ("GMM",           "Gaussian Mixture Model — soft probabilistic clustering"),
        ("Grad-CAM",      "Gradient-weighted Class Activation Mapping — visual explanations"),
        ("HbA1c",         "Glycated Hemoglobin — diabetes control marker"),
        ("ICD-10",        "International Classification of Diseases, 10th Revision"),
        ("IQR",           "Interquartile Range — used for outlier detection"),
        ("LOS",           "Length of Stay"),
        ("LSTM",          "Long Short-Term Memory — recurrent neural network"),
        ("MCI",           "Mass Casualty Incident"),
        ("MEWS",          "Modified Early Warning Score — validated deterioration score"),
        ("NER",           "Named Entity Recognition"),
        ("OSRM",          "Open Source Routing Machine — road-network routing"),
        ("PCA",           "Principal Component Analysis — linear dim reduction"),
        ("PHI",           "Protected Health Information (HIPAA)"),
        ("RAG",           "Retrieval-Augmented Generation"),
        ("SHAP",          "SHapley Additive exPlanations — game-theoretic feature importance"),
        ("SpO2",          "Peripheral oxygen saturation"),
        ("t-SNE",         "t-Distributed Stochastic Neighbor Embedding"),
        ("UMAP",          "Uniform Manifold Approximation and Projection"),
        ("WBC",           "White Blood Cell count"),
    ]
    # 2 columns
    for i, (term, defn) in enumerate(terms):
        col = i % 2
        row = i // 2
        cx = 0.4 + col * 6.3
        cy = y0+0.15 + row * 0.5
        add_text(s, term, cx, cy+0.05, 1.2, 0.3, size=11, bold=True, color=INDIGO, font="Consolas")
        add_text(s, defn, cx+1.3, cy+0.05, 4.8, 0.3, size=10, color=SLATE)

# ═════════════════════════════════════════════════════════════════════════════
# SECTION 9 — IMPACT + CLOSE
# ═════════════════════════════════════════════════════════════════════════════
def slide_impact():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    y0 = header(s, "Expected Impact", "What a 500-bed hospital deployment would actually deliver")
    # stats
    stats = [
        ("15→2 min",  "Triage time per patient",   "Automated risk stratification", INDIGO),
        ("↓ ICU",     "Preventable transfers",     "MEWS-based early detection",   VIOLET),
        ("100%",      "Drug safety check",         "Every med list auto-scanned",  RED),
        ("↓ 30%",     "Radiologist routine load",  "AI X-ray pre-reading",         TEAL),
    ]
    for i, (big, small, desc, c) in enumerate(stats):
        cx = 0.4 + i * 3.2
        add_rect(s, cx, y0+0.1, 3.0, 2.0, c)
        add_text(s, big, cx, y0+0.25, 3.0, 0.7, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, small, cx, y0+0.95, 3.0, 0.35, size=12, color=LIGHT, align=PP_ALIGN.CENTER, italic=True)
        add_text(s, desc, cx+0.15, y0+1.4, 2.7, 0.5, size=10, color=LIGHT, align=PP_ALIGN.CENTER)
    # key innovations
    add_text(s, "FIVE KEY INNOVATIONS", 0.4, y0+2.4, 12.5, 0.35, size=14, bold=True, color=INDIGO)
    innovations = [
        ("Zero-label risk stratification",
         "Meaningful risk tiers from raw patient data — no pre-labelled training data needed"),
        ("Graceful degradation architecture",
         "Every component has a fallback: SHAP → zero stubs, LSTM → linear, SpaCy → regex, FAISS → exact"),
        ("Field alias resolution",
         "Handles diverse EHR column naming conventions — works with any hospital data export"),
        ("Clinical transparency",
         "Every prediction includes SHAP explanations, confidence scores, natural-language care plans"),
        ("Integrated emergency ecosystem",
         "Risk scoring → ambulance dispatch → MCI board — covers full emergency care workflow"),
    ]
    for i, (t, d) in enumerate(innovations):
        col = i % 3
        row = i // 3
        cx = 0.4 + col * 4.3
        cy = y0+2.85 + row*1.6
        add_round(s, cx, cy, 4.1, 1.45, WHITE, line=BORDER)
        add_rect(s, cx, cy, 0.12, 1.45, INDIGO)
        add_text(s, f"{i+1}. {t}", cx+0.25, cy+0.1, 3.8, 0.35, size=12, bold=True, color=INDIGO)
        add_text(s, d, cx+0.25, cy+0.5, 3.8, 0.9, size=10, color=SLATE)

def slide_summary():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, 13.333, 7.5, INDIGO)
    add_rect(s, 0, 5.0, 13.333, 2.5, VIOLET)
    add_text(s, "MediCluster", 0, 0.8, 13.333, 1.2, size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "From raw patient data to clinical action — in one platform.",
             0, 2.0, 13.333, 0.5, size=18, color=LIGHT, align=PP_ALIGN.CENTER, italic=True)
    # 5 stat columns
    stats = [
        ("4",  "Clustering\nalgorithms"),
        ("5",  "Imaging\nmodels"),
        ("20+","ML/NLP\nendpoints"),
        ("13", "Frontend\npages"),
        ("3",  "AI-powered\ndashboards"),
    ]
    for i, (n, l) in enumerate(stats):
        cx = 0.8 + i * 2.4
        add_rect(s, cx, 2.9, 2.1, 1.6, INDIGO_D)
        add_text(s, n, cx, 3.05, 2.1, 0.7, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, l, cx, 3.8, 2.1, 0.7, size=11, color=LIGHT, align=PP_ALIGN.CENTER)
    add_text(s, "ML Clustering  ·  LungLens Imaging  ·  CliniQ NLP  ·  ARIA Dispatch",
             0, 5.3, 13.333, 0.5, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Built for healthcare professionals who demand precision, transparency, and zero-fail operation.",
             0, 5.9, 13.333, 0.5, size=12, color=LIGHT, align=PP_ALIGN.CENTER, italic=True)
    add_text(s, "CtrlAltElite  ·  M S Engineering College  ·  May 2026  ·  http://localhost:5173",
             0, 6.6, 13.333, 0.4, size=11, color=RGBColor(0xC7, 0xD2, 0xFE), align=PP_ALIGN.CENTER, italic=True)
    add_text(s, "Thank you", 0, 4.6, 13.333, 0.4, size=14, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.CENTER, italic=True)

# ═════════════════════════════════════════════════════════════════════════════
# BUILD
# ═════════════════════════════════════════════════════════════════════════════
title_slide()
agenda()
slide_problem()
slide_goals_users()
slide_architecture()
slide_tech_stack()
slide_data_flow()
slide_preprocessing()
slide_clustering_algos()
slide_risk_scoring()
slide_shap()
slide_anomalies()
slide_automl()
slide_nlp()
slide_forecasting()
slide_imaging()
slide_rag_drug()
slide_pages()
slide_dispatch()
slide_api()
slide_deployment()
slide_performance()
slide_security()
slide_testing()
slide_roadmap()
slide_limitations()
slide_glossary()
slide_impact()
slide_summary()

out = r"C:\Users\Admin\Desktop\Patient Health Risk Segregation\MediCluster_Definitive.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
