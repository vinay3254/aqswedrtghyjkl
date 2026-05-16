from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import datetime

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

INDIGO  = RGBColor(79,  70,  229)
VIOLET  = RGBColor(124, 58,  237)
SLATE   = RGBColor(51,  65,  85)
MUTED   = RGBColor(100, 116, 139)
WHITE   = RGBColor(255, 255, 255)
LIGHT   = RGBColor(238, 242, 255)
RED     = RGBColor(220, 38,  38)
GREEN   = RGBColor(22,  163, 74)
AMBER   = RGBColor(217, 119, 6)
BG      = RGBColor(248, 250, 252)

BLANK = prs.slide_layouts[6]   # completely blank

def solid_fill(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, color, radius=False):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    solid_fill(shape, color)
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=SLATE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.15, INDIGO)
    add_text(slide, title, 0.4, 0.1, 10, 0.65, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.72, 10, 0.38, size=13, color=RGBColor(199,210,254))

def bullet_box(slide, items, x, y, w, h, size=12, color=SLATE, title=None, title_color=INDIGO):
    if title:
        add_text(slide, title, x, y, w, 0.35, size=13, bold=True, color=title_color)
        y += 0.38
        h -= 0.38
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = u"•  " + item
        run.font.size = Pt(size)
        run.font.color.rgb = color

def card(slide, x, y, w, h, header_text, body_lines, hdr_color=INDIGO, body_size=11):
    add_rect(slide, x, y, w, 0.42, hdr_color)
    add_text(slide, header_text, x+0.1, y+0.04, w-0.2, 0.35, size=12, bold=True, color=WHITE)
    add_rect(slide, x, y+0.42, w, h-0.42, WHITE)
    body_box = slide.shapes.add_textbox(Inches(x+0.12), Inches(y+0.5), Inches(w-0.24), Inches(h-0.58))
    tf = body_box.text_frame
    tf.word_wrap = True
    first = True
    for line in body_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = u"•  " + line
        run.font.size = Pt(body_size)
        run.font.color.rgb = SLATE

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, INDIGO)
add_rect(s, 0, 4.8, 13.33, 2.7, VIOLET)

add_text(s, "MediCluster", 1.0, 1.2, 11, 1.5, size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Patient Health Risk Segregation Platform",
         1.0, 2.9, 11, 0.7, size=22, color=RGBColor(199,210,254), align=PP_ALIGN.CENTER)
add_text(s, "AI-Powered Clinical Decision Support  |  " + datetime.date.today().strftime("%B %Y"),
         1.0, 3.7, 11, 0.5, size=14, color=RGBColor(165,180,252), align=PP_ALIGN.CENTER)

for i, label in enumerate(["ML Clustering", "LungLens Imaging", "CliniQ NLP", "ARIA Dispatch"]):
    cx = 1.0 + i * 2.85
    add_rect(s, cx, 5.1, 2.5, 1.6, RGBColor(55,48,163))
    add_text(s, label, cx+0.1, 5.4, 2.3, 0.8, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT IS MEDICLUSTER
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "What is MediCluster?", "A full-stack AI clinical decision support platform")

add_text(s, "MediCluster combines machine learning, deep learning, NLP, and AI to help hospitals "
            "identify high-risk patients, analyze medical images, manage emergencies, and support clinical decisions.",
         0.5, 1.3, 12.3, 0.8, size=13, color=MUTED)

stats = [
    ("4", "Clustering\nAlgorithms"),
    ("6", "Imaging\nModels"),
    ("5", "Languages\nSupported"),
    ("14+", "ML/NLP\nEndpoints"),
    ("3", "AI-Powered\nDashboards"),
]
for i, (num, lbl) in enumerate(stats):
    cx = 0.5 + i * 2.5
    add_rect(s, cx, 2.3, 2.1, 1.8, INDIGO)
    add_text(s, num,  cx, 2.4, 2.1, 0.9, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, lbl,  cx, 3.25, 2.1, 0.7, size=11, color=RGBColor(199,210,254), align=PP_ALIGN.CENTER)

bullet_box(s, [
    "React 18 + Vite + Tailwind CSS frontend (port 3000)",
    "Node.js 25 + Express API gateway with Anthropic SDK (port 5000)",
    "Python 3.11 + Flask ML engine with PyTorch + scikit-learn (port 8000)",
    "MongoDB 7 with GridFS for images and DICOM files",
], 0.5, 4.3, 12.3, 2.8, size=12)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "System Architecture", "Three-tier microservice design")

layers = [
    (INDIGO,                 "Browser  (React + Vite)",    "port 3000",  "Dashboard, Imaging, CliniQ, Dispatch, Triage, Reminders"),
    (RGBColor(37,99,235),    "API Gateway  (Node.js)",      "port 5000",  "REST routing, MongoDB persistence, Claude API proxy"),
    (RGBColor(5,150,105),    "ML Engine  (Python Flask)",   "port 8000",  "Clustering, NLP, DenseNet inference, forecasting"),
    (RGBColor(71,85,105),    "MongoDB",                     "port 27017", "Patient data, cluster results, GridFS media storage"),
]
for i, (clr, name, port, desc) in enumerate(layers):
    y = 1.35 + i * 1.38
    add_rect(s, 0.4, y, 12.5, 1.15, clr)
    add_text(s, name, 0.6, y+0.08, 3.5, 0.5, size=14, bold=True, color=WHITE)
    add_text(s, port, 4.3, y+0.08, 1.5, 0.5, size=13, color=RGBColor(199,210,254), italic=True)
    add_text(s, desc, 5.9, y+0.08, 6.8, 0.9, size=12, color=RGBColor(224,231,255))
    if i < 3:
        add_text(s, u"↓", 6.0, y+1.0, 1.5, 0.35, size=18, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ML CLUSTERING
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Patient Risk Clustering", "Core ML feature — unsupervised risk segmentation from CSV vitals")

algos = [
    (INDIGO,               "K-Means",      ["Spherical clusters", "Fast convergence", "Auto-K via Elbow+Silhouette"]),
    (RGBColor(37,99,235),  "DBSCAN",       ["Arbitrary shape clusters", "Noise/outlier isolation", "No K needed"]),
    (RGBColor(126,34,206), "Hierarchical", ["Dendrogram visualization", "Ward/complete linkage", "Any cluster count"]),
    (RGBColor(5,150,105),  "GMM",          ["Soft membership", "Probabilistic assignment", "Covariance control"]),
]
for i, (clr, name, pts) in enumerate(algos):
    cx = 0.4 + i * 3.2
    card(s, cx, 1.3, 3.0, 3.0, name, pts, hdr_color=clr)

bullet_box(s, [
    "AutoML: automatic K selection (Elbow + Silhouette sweep K=2..10)",
    "SHAP explainability: per-patient feature contribution scores shown per cluster",
    "Isolation Forest anomaly detection flags statistical outliers within clusters",
    "UMAP + t-SNE 2D projections as alternatives to PCA scatter plot",
    "Risk tiers assigned: Critical / High / Moderate / Low / Noise",
], 0.4, 4.5, 12.5, 2.6, size=12)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — LUNGLENS + IMAGING
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Medical Imaging — LungLens + DenseNet", "Six model options for chest X-ray and general medical image analysis")

models = [
    "CheXNet (CheXpert) — recommended",
    "DenseNet121 All datasets",
    "ResNet50 High-res 512px",
    "DenseNet121 NIH ChestX-ray14",
    "DenseNet121 PadChest",
    "LungLens (Claude Vision AI)",
]
for i, m in enumerate(models):
    col = i % 3
    row = i // 3
    cx = 0.4 + col * 4.3
    cy = 1.35 + row * 1.45
    clr = VIOLET if "LungLens" in m else INDIGO
    add_rect(s, cx, cy, 4.0, 1.2, clr)
    add_text(s, m, cx+0.15, cy+0.3, 3.7, 0.65, size=12, bold=True, color=WHITE)

bullet_box(s, [
    "Each finding returns: label, confidence %, severity badge (HIGH / MODERATE / LOW), cause, medications, prevention",
    "LungLens uses Claude Sonnet 4.6 Vision — works on any medical image, not just chest X-rays",
    "Grad-CAM heatmaps show which image region triggered each pathology detection",
    "AI Deep Explanation: Claude generates plain-English paragraph summary of all findings",
    "DICOM (.dcm) file support with CLAHE contrast enhancement preprocessing",
], 0.4, 4.35, 12.5, 2.7, size=12)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — CLINIQ MULTILINGUAL
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "CliniQ — Multilingual AI Assistant", "Vision-capable AI chat for medical image Q&A in 5 Indian + English languages")

langs = [
    (INDIGO,               "English",  "Latin script",      "en"),
    (RGBColor(220,38,38),  "Hindi",    "Devanagari",        "hi"),
    (RGBColor(5,150,105),  "Kannada",  "Kannada script",    "kn"),
    (RGBColor(217,119,6),  "Telugu",   "Telugu script",     "te"),
    (RGBColor(126,34,206), "Tamil",    "Tamil script",      "ta"),
]
for i, (clr, name, script, code) in enumerate(langs):
    cx = 0.4 + i * 2.5
    add_rect(s, cx, 1.35, 2.2, 1.8, clr)
    add_text(s, name,   cx, 1.45, 2.2, 0.65, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, script, cx, 2.1,  2.2, 0.5,  size=11, color=RGBColor(224,231,255), align=PP_ALIGN.CENTER)
    add_text(s, code,   cx, 2.6,  2.2, 0.4,  size=13, bold=True, color=RGBColor(199,210,254), align=PP_ALIGN.CENTER)

bullet_box(s, [
    "Language selector dropdown in CliniQ header — always visible, defaults to English",
    "Claude Sonnet 4.6 receives a system prompt instructing it to respond entirely in the selected language",
    "All section headings (Condition Name, Prevention, Treatment, Medications) translated automatically",
    "Medical disclaimer translated — 'Educational purposes only, consult a physician'",
    "Multi-turn conversation maintains language across follow-up questions",
    "Resets to English on 'New Chat'",
], 0.4, 3.35, 12.5, 3.7, size=12)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — NLP FEATURES
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Clinical NLP Pipeline", "Natural language processing for clinical notes, prescriptions, and drug safety")

nlp_cards = [
    (INDIGO,               "NER + ICD-10",        ["Named entity recognition", "Drugs, diagnoses, symptoms", "ICD-10 code suggestion"]),
    (RGBColor(220,38,38),  "Drug Interactions",   ["12-pair knowledge base", "Major / moderate severity", "Warfarin, metformin, etc."]),
    (RGBColor(5,150,105),  "Prescription Extract",["Free-text prescription input", "Extract drug, dose, frequency", "Claude fallback if regex fails"]),
    (RGBColor(126,34,206), "Patient Trajectory",  ["Classify note sentiment", "Improving / stable / deteriorating", "Risk keyword lexicon"]),
]
for i, (clr, name, pts) in enumerate(nlp_cards):
    cx = 0.4 + i * 3.2
    card(s, cx, 1.35, 3.0, 2.7, name, pts, hdr_color=clr)

bullet_box(s, [
    "Vital Signs Forecasting: LSTM + Prophet time-series for HR, BP, SpO2, temperature, respiratory rate",
    "MEWS Score: Modified Early Warning Score — returns alert level: low / moderate / high / critical",
    "RAG Chatbot: FAISS vector store + TF-IDF fallback for patient Q&A on clinical notes",
    "Anomaly Detection: Isolation Forest flags statistical outliers across patient features",
], 0.4, 4.25, 12.5, 2.8, size=12)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ARIA DISPATCH + TRIAGE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "ARIA Dispatch + MCI Triage Board", "AI-powered emergency response and mass casualty incident management")

add_rect(s, 0.4, 1.35, 6.0, 5.7, WHITE)
add_rect(s, 0.4, 1.35, 6.0, 0.45, INDIGO)
add_text(s, "ARIA — Ambulance Dispatch", 0.55, 1.38, 5.7, 0.38, size=13, bold=True, color=WHITE)
bullet_box(s, [
    "Live Leaflet.js dispatcher map with ambulance positions",
    "Real-time status: available / en_route / at_scene",
    "Claude selects nearest ambulance + best hospital",
    "Hospital match by bed count and specialty",
    "Route polyline visualization on map",
    "Driver dashboard with patient alert overlays",
    "Full dispatch history (resolve/cancel)",
], 0.55, 1.92, 5.7, 4.8, size=11)

add_rect(s, 6.9, 1.35, 6.0, 5.7, WHITE)
add_rect(s, 6.9, 1.35, 6.0, 0.45, VIOLET)
add_text(s, "MCI Triage Board", 7.05, 1.38, 5.7, 0.38, size=13, bold=True, color=WHITE)
bullet_box(s, [
    "Voice-to-triage: speak/type patient condition",
    "Claude extracts vitals + assigns P1/P2/P3/P4",
    "Stateless scoring: instant score, no DB write",
    "Hospital allocation by triage category",
    "Pathway flags (trauma/cardiac/respiratory)",
    "Live patient cards with vitals history",
    "Transport status tracking per patient",
], 7.05, 1.92, 5.7, 4.8, size=11)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — TECH STACK
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Technology Stack", "Full-stack AI platform built on open-source and cloud AI")

stack = [
    ("Frontend",      INDIGO,               ["React 18", "Vite 5", "Tailwind CSS 3", "Recharts + D3.js", "React Router 6"]),
    ("Backend",       RGBColor(37,99,235),  ["Node.js 25", "Express 4", "Mongoose 8", "Multer", "Anthropic SDK 0.95"]),
    ("ML / AI",       RGBColor(126,34,206), ["Python 3.11", "Flask 3", "scikit-learn 1.4", "PyTorch 2.2", "torchxrayvision 1.0"]),
    ("NLP",           RGBColor(5,150,105),  ["spaCy 3.7", "scispaCy 0.5", "HuggingFace 4.40", "FAISS 1.8", "SHAP 0.45"]),
    ("Data / AI",     RGBColor(217,119,6),  ["MongoDB 7", "GridFS", "UMAP-learn", "Prophet 1.1", "Claude Sonnet 4.6"]),
]
for i, (title, clr, items) in enumerate(stack):
    cx = 0.4 + i * 2.5
    card(s, cx, 1.35, 2.3, 5.7, title, items, hdr_color=clr, body_size=11)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — RECENT ADDITIONS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Recent Feature Additions", "New capabilities added in latest development cycle")

updates = [
    (INDIGO,               "Multilingual CliniQ",
     ["Language selector: EN / HI / KN / TE / TA",
      "Claude responds in full selected language",
      "All headings + disclaimer translated",
      "Resets to English on New Chat"]),
    (VIOLET,               "LungLens (Claude Vision)",
     ["6th imaging model option",
      "Works on any medical image type",
      "Same structured output as DenseNet",
      "Confidence, severity, medications, prevention"]),
    (RGBColor(5,150,105),  "Shared Anthropic Client",
     ["Single SDK client with 120s timeout",
      "Fixed APIConnectionTimeoutError",
      "Re-used across all 4 Claude endpoints",
      "Faster due to connection reuse"]),
    (RGBColor(217,119,6),  "Dispatch Maps",
     ["Live Leaflet.js dispatcher map",
      "Driver dashboard with patient alerts",
      "Route polyline visualization",
      "Ambulance + hospital overlays"]),
]
for i, (clr, title, pts) in enumerate(updates):
    cx = 0.4 + i * 3.2
    card(s, cx, 1.35, 3.0, 5.7, title, pts, hdr_color=clr)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — THANK YOU
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, INDIGO)
add_rect(s, 0, 5.0, 13.33, 2.5, VIOLET)
add_text(s, "Thank You", 0, 1.5, 13.33, 1.8, size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "MediCluster — Patient Health Risk Segregation Platform",
         0, 3.3, 13.33, 0.7, size=18, color=RGBColor(199,210,254), align=PP_ALIGN.CENTER)
add_text(s, "http://localhost:3000",
         0, 4.1, 13.33, 0.5, size=13, color=RGBColor(165,180,252), align=PP_ALIGN.CENTER, italic=True)

for i, label in enumerate(["ML Clustering", "LungLens AI", "CliniQ NLP", "ARIA Dispatch"]):
    cx = 1.2 + i * 2.85
    add_rect(s, cx, 5.4, 2.5, 1.4, RGBColor(55,48,163))
    add_text(s, label, cx, 5.75, 2.5, 0.65, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

out = r"c:\Users\Admin\Desktop\Patient Health Risk Segregation\MediCluster_Presentation.pptx"
prs.save(out)
print("Saved:", out)
